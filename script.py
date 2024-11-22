from pptx import Presentation
import json
import os

# Carregar o JSON
with open("dados.json", "r", encoding="utf-8") as arquivo:
    participantes = json.load(arquivo)

# Criar pasta para salvar os certificados gerados
os.makedirs("certificados", exist_ok=True)

# Gerar certificados
for participante in participantes:
    nome = participante.get("nome")
    cpf = participante.get("CPF  ")  # Note o espaço extra aqui, caso ainda esteja presente

    if not nome or not cpf:
        print(f"Erro: Dados ausentes para {participante}")
        continue

    certificado = Presentation("modelo.pptx")
    for slide in certificado.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:  # Verifica se o shape possui texto
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:  # Itera sobre cada segmento de texto
                        run.text = run.text.replace("{{NOME}}", nome)
                        run.text = run.text.replace("{{CPF}}", cpf)

    nome_arquivo = f"certificados/{nome}_certificado.pptx"
    certificado.save(nome_arquivo)

print("Certificados gerados com sucesso!")
